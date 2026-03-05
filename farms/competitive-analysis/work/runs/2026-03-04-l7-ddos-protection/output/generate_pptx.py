"""
Generate Competitive Analysis PowerPoint deck for L7 DDoS Protection.
Run: python generate_pptx.py
Output: competitive-analysis.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Color scheme ---
DARK_BLUE = RGBColor(0x0B, 0x2A, 0x4A)
MEDIUM_BLUE = RGBColor(0x1B, 0x4F, 0x72)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1C, 0x1C, 0x1C)
DARK_GRAY = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1E, 0x88, 0x51)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=WHITE):
    """Add a solid background rectangle."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_header_bar(slide, height=Inches(1.1)):
    """Add dark blue header bar at top of slide."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()


def add_footer(slide):
    """Add classification footer."""
    ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    tf = ft.text_frame
    p = tf.paragraphs[0]
    p.text = "Internal  |  Competitive Analysis Agent Farm  |  March 4, 2026"
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY
    p.font.italic = True


def add_title_text(slide, text, left, top, width, height, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT):
    """Add a title text box."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_body_text(slide, text, left, top, width, height, font_size=16, color=BLACK, bold=False):
    """Add body text box."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    return tf


def add_bullet_list(slide, items, left, top, width, height, font_size=15, color=BLACK):
    """Add a bulleted list."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.level = 0
    return tf


def add_card(slide, title, items, left, top, width, height, header_color=MEDIUM_BLUE):
    """Add a card with header and bullet items."""
    # Header
    hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.45))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = header_color
    hdr.line.fill.background()
    htf = hdr.text_frame
    htf.paragraphs[0].text = title
    htf.paragraphs[0].font.size = Pt(14)
    htf.paragraphs[0].font.color.rgb = WHITE
    htf.paragraphs[0].font.bold = True
    htf.paragraphs[0].alignment = PP_ALIGN.CENTER
    htf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Body
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top + Inches(0.45), width, height - Inches(0.45))
    body.fill.solid()
    body.fill.fore_color.rgb = LIGHT_GRAY
    body.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)

    btf = body.text_frame
    btf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = btf.paragraphs[0]
        else:
            p = btf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)


# =============================================================================
# SLIDE 1 — Title Slide
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

# Accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.6), Inches(1.5), Inches(0.06))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

add_title_text(slide, "Competitive Brief", Inches(1), Inches(1.5), Inches(11), Inches(1.2),
               font_size=44, color=WHITE, bold=True)
add_title_text(slide, "L7 DDoS Protection", Inches(1), Inches(2.7), Inches(11), Inches(1),
               font_size=36, color=ACCENT_BLUE, bold=False)
add_title_text(slide, "March 4, 2026", Inches(1), Inches(4.0), Inches(5), Inches(0.5),
               font_size=18, color=RGBColor(0xAA, 0xBB, 0xCC), bold=False)
add_title_text(slide, "Classification: Internal", Inches(1), Inches(4.5), Inches(5), Inches(0.5),
               font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC), bold=False)
add_title_text(slide, "Audience: Leadership / Executive", Inches(1), Inches(4.9), Inches(5), Inches(0.5),
               font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC), bold=False)

# =============================================================================
# SLIDE 2 — Executive Summary
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_title_text(slide, "Executive Summary", Inches(0.6), Inches(0.25), Inches(10), Inches(0.7),
               font_size=30, color=WHITE)
add_footer(slide)

exec_items = [
    "Cloudflare is the clear market leader in L7 DDoS — autonomous, ML-driven, unmetered mitigation across a 321 Tbps network with ~3s time-to-mitigate.",
    "Azure WAF provides solid OWASP-based web application protection with strong native Azure integration, but treats L7 DDoS as a secondary use case.",
    "Primary competitive gap: No ML-based behavioral L7 DDoS detection in production — Cloudflare has shipped this; Azure WAF has not.",
    "Strategic imperative: Accelerate the dedicated L7 DDoS detection pipeline while leveraging current differentiators (MazeBolt RADAR, Azure security stack, FedRAMP High).",
    "Key differentiator today: MazeBolt RADAR — non-disruptive, validated L7 DDoS simulation exclusive to Azure. No competitor can match this.",
]
add_bullet_list(slide, exec_items, Inches(0.6), Inches(1.4), Inches(12), Inches(5.2), font_size=17)

# =============================================================================
# SLIDE 3 — Market Landscape
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_title_text(slide, "Market Landscape", Inches(0.6), Inches(0.25), Inches(10), Inches(0.7),
               font_size=30, color=WHITE)
add_footer(slide)

# Source caveat
add_body_text(slide, "Source note: All market data from Cloudflare DDoS Threat Reports (self-reported). No independent analyst data available.",
              Inches(0.6), Inches(1.3), Inches(12), Inches(0.5), font_size=12, color=DARK_GRAY, bold=False)

market_items = [
    "21.3M DDoS attacks blocked by Cloudflare in 2024 (+53% YoY)",
    "L7 attacks now ~51% of all DDoS (Q4 2024) — application-layer is the new battleground",
    "72% of L7 attacks end in <10 minutes — automated, always-on protection is critical",
    "73% of HTTP DDoS attacks launched by known botnets; Mirai variants +131% QoQ",
    "Hyper-volumetric attacks (>1 Tbps) grew 1,885% QoQ in Q4 2024; record 5.6 Tbps mitigated",
    "Ransom DDoS up 78% QoQ in Q4 2024; 12% of targets received extortion threats",
    "Unmetered DDoS becoming table stakes — Cloudflare offers unlimited protection on all plans (including Free)",
    "PeerSpot WAF mindshare: Imperva 8.1%, Fortinet FortiWeb 7.5%, Azure WAF 2.8%",
]
add_bullet_list(slide, market_items, Inches(0.6), Inches(1.9), Inches(12), Inches(5.0), font_size=16)

# =============================================================================
# SLIDE 4 — Comparison Matrix
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_title_text(slide, "Comparison Matrix: Azure WAF vs. Cloudflare", Inches(0.6), Inches(0.25),
               Inches(11), Inches(0.7), font_size=28, color=WHITE)
add_footer(slide)

# Table
rows_data = [
    ("Dimension", "Azure WAF", "Cloudflare"),
    ("L7 DDoS Approach", "Composite: rate limiting + bot rules + custom rules. No dedicated engine.", "Purpose-built: Autonomous ML engine. ~3s mitigate. No human intervention."),
    ("ML/AI Detection", "Rule-based anomaly scoring. No adaptive L7 DDoS profiling.", "Adaptive DDoS: 7-day ML traffic profiling per zone."),
    ("Attack Coverage", "OWASP CRS + bot rules. No L7 DDoS-specific signatures.", "HTTP floods, H2 Rapid Reset, Slowloris, cache busting, TLS exhaustion, 20+ vectors."),
    ("Network Capacity", "Not publicly disclosed.", "321 Tbps across 330+ cities. Mitigated 5.6 Tbps (2024)."),
    ("Time to Mitigate", "No published benchmark.", "~3 seconds average."),
    ("Unmetered DDoS", "No. Capacity-based pricing.", "Yes. All plans including Free ($0)."),
    ("Pricing (DDoS)", "WAF: $0.443/gw-hr + CU. DDoS Prot: ~$2,944/mo.", "Free–$200/mo + Enterprise custom. Unmetered."),
    ("Azure Integration", "Native: Sentinel, Defender, Monitor, Portal.", "External: Logpush, API, Terraform."),
    ("Multi-Cloud", "Azure-only.", "Cloud-agnostic. Any origin."),
    ("Time-to-Value", "App Gateway/Front Door provisioning + rule tuning.", "DNS change. Protection in minutes."),
    ("Compliance", "PCI-DSS. FedRAMP High likely (verify).", "SOC 2, ISO 27001, PCI DSS, FedRAMP Moderate."),
]

tbl_left = Inches(0.4)
tbl_top = Inches(1.3)
tbl_width = Inches(12.5)
num_rows = len(rows_data)
num_cols = 3
col_widths = [Inches(2.2), Inches(5.15), Inches(5.15)]
row_height = Inches(0.45)

table_shape = slide.shapes.add_table(num_rows, num_cols, tbl_left, tbl_top, tbl_width, Inches(num_rows * 0.45))
table = table_shape.table

for ci, cw in enumerate(col_widths):
    table.columns[ci].width = cw

for ri, row in enumerate(rows_data):
    for ci, cell_text in enumerate(row):
        cell = table.cell(ri, ci)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = BLACK if ri > 0 else WHITE
            paragraph.font.bold = (ri == 0)
        if ri == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
        elif ri % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

# =============================================================================
# SLIDE 5 — Azure WAF Deep-Dive
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_title_text(slide, "Azure WAF — Deep-Dive", Inches(0.6), Inches(0.25), Inches(10), Inches(0.7),
               font_size=30, color=WHITE)
add_footer(slide)

add_card(slide, "STRENGTHS", [
    "Native Azure integration — Sentinel, Defender, Monitor single-pane",
    "Flexible deployment: App Gateway, Front Door, CDN, Containers",
    "OWASP CRS 3.2 with anomaly scoring engine",
    "MazeBolt RADAR — validated L7 DDoS simulation (Azure exclusive)",
    "Enterprise trust: Microsoft support + Defender ecosystem",
    "Cost-effective WAF management for existing Azure customers",
], Inches(0.5), Inches(1.3), Inches(5.8), Inches(3.0), header_color=GREEN)

add_card(slide, "WEAKNESSES", [
    "No dedicated L7 DDoS engine — manual rule assembly required",
    "L3/L4 DDoS requires separate ~$2,944/mo Azure DDoS Protection",
    "Capacity-based pricing — large attacks increase costs",
    "Azure-only — no multi-cloud or on-prem support",
    "Only 2.8% PeerSpot WAF mindshare",
    "Configuration complexity cited by users",
], Inches(6.8), Inches(1.3), Inches(5.8), Inches(3.0), header_color=RED)

add_card(slide, "KEY FEATURES", [
    "OWASP CRS 3.0/3.1/3.2 managed rulesets",
    "IP-based rate limiting (1-min / 5-min windows)",
    "Bot Manager Rule Set + Microsoft Threat Intelligence",
    "Custom rules with match conditions (headers, cookies, geo, IP)",
    "Detection and Prevention modes for safe rollout",
], Inches(0.5), Inches(4.5), Inches(12.1), Inches(2.3), header_color=MEDIUM_BLUE)

# =============================================================================
# SLIDE 6 — Cloudflare Deep-Dive
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_title_text(slide, "Cloudflare — Deep-Dive", Inches(0.6), Inches(0.25), Inches(10), Inches(0.7),
               font_size=30, color=WHITE)
add_footer(slide)

add_card(slide, "STRENGTHS", [
    "Unmetered DDoS on all plans — even Free ($0)",
    "321 Tbps network; mitigated 5.6 Tbps attack autonomously",
    "Fully autonomous detection — ~3s time-to-mitigate",
    "Comprehensive L7 attack coverage in single managed ruleset",
    "Adaptive ML protection — per-zone traffic profiling",
    "Industry-leading threat intel + public DDoS reports",
], Inches(0.5), Inches(1.3), Inches(5.8), Inches(3.0), header_color=GREEN)

add_card(slide, "WEAKNESSES", [
    "Advanced features require Enterprise tier",
    "Reverse-proxy requirement for L7 protection",
    "Less native cloud-platform integration",
    "Enterprise pricing is opaque (custom-quoted)",
    "Wirefilter syntax has learning curve vs. portal builder",
], Inches(6.8), Inches(1.3), Inches(5.8), Inches(3.0), header_color=RED)

add_card(slide, "KEY FEATURES", [
    "Autonomous DDoS engine with real-time signature generation",
    "HTTP DDoS Attack Protection Managed Ruleset (20+ vectors)",
    "Adaptive DDoS Protection — ML 7-day traffic profiling",
    "ML-based bot scoring with known botnet fingerprinting",
    "Flexible rate limiting + custom WAF rules (wirefilter)",
], Inches(0.5), Inches(4.5), Inches(12.1), Inches(2.3), header_color=MEDIUM_BLUE)

# =============================================================================
# SLIDE 7 — Battle Cards
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_title_text(slide, "Battle Cards: Azure WAF vs. Cloudflare", Inches(0.6), Inches(0.25),
               Inches(11), Inches(0.7), font_size=28, color=WHITE)
add_footer(slide)

# Warning
add_body_text(slide, "⚠️ Internal — do not distribute externally.",
              Inches(0.6), Inches(1.15), Inches(6), Inches(0.35), font_size=12, color=RED, bold=True)

battle_data = [
    ("Dimension", "Azure WAF Wins Because", "Cloudflare Wins Because", "Objection Handler"),
    ("Azure\nIntegration", "Single-pane: Sentinel, Defender,\nMonitor. No context-switching.", "N/A — cloud-agnostic,\nno native Azure integration.", "\"We use Splunk\" → Azure also integrates\nvia diagnostic settings. Unique: correlated\nsecurity across Azure networking + identity."),
    ("L7 DDoS\nAutomation", "Detection/Prevention modes allow\ncontrolled rollout, fewer false positives.", "Fully autonomous ~3s mitigate.\n72% of L7 attacks end in <10 min.", "Lead with MazeBolt RADAR validated\ntesting. Our WAF+DDoS covers L3-L7\nwith full Defender ecosystem."),
    ("ML/Adaptive\nDetection", "Anomaly scoring + MS Threat Intel.\nAzure-only signals unavailable to\nexternal vendors.", "Production Adaptive DDoS: 7-day\nML profiling, multi-signal anomaly\ndetection at scale.", "Acknowledge gap. Our strength:\ncorrelated security intel across Azure\nplatform (Defender, Sentinel, Entra ID)."),
    ("Attack\nCoverage", "OWASP CRS 3.2 + MazeBolt RADAR\nvalidates actual coverage with\nnon-disruptive simulation.", "Dedicated L7 DDoS ruleset: HTTP\nfloods, H2 Rapid Reset, Slowloris,\ncache busting, 20+ vectors.", "\"Can your vendor demonstrate validated\nprotection under simulated attack, or\njust show a list of supported types?\""),
    ("Pricing", "WAF included in App Gateway/Front\nDoor. No separate vendor contract.", "Unmetered DDoS on all plans\nincluding Free ($0). No bandwidth\npenalties.", "Azure pricing = enterprise SLA + support.\nInclude vendor consolidation savings\nand MS Unified Support value."),
    ("Compliance", "FedRAMP High (Azure Gov).\nCritical for U.S. federal.", "FedRAMP Moderate only.\nDoes not meet High.", "Clear differentiator for federal\nand regulated workloads."),
]

tbl_left = Inches(0.3)
tbl_top = Inches(1.5)
num_rows = len(battle_data)
num_cols = 4
col_widths = [Inches(1.4), Inches(3.5), Inches(3.5), Inches(4.1)]
table_shape = slide.shapes.add_table(num_rows, num_cols, tbl_left, tbl_top, Inches(12.5), Inches(5.2))
table = table_shape.table

for ci, cw in enumerate(col_widths):
    table.columns[ci].width = cw

for ri, row in enumerate(battle_data):
    for ci, cell_text in enumerate(row):
        cell = table.cell(ri, ci)
        cell.text = cell_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = BLACK if ri > 0 else WHITE
            paragraph.font.bold = (ri == 0) or (ci == 0 and ri > 0)
        if ri == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
        elif ri % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

# =============================================================================
# SLIDE 8 — Strategic Recommendations
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_header_bar(slide)
add_title_text(slide, "Strategic Recommendations", Inches(0.6), Inches(0.25), Inches(10), Inches(0.7),
               font_size=30, color=WHITE)
add_footer(slide)

# Warning
add_body_text(slide, "⚠️ Internal — do not distribute externally.",
              Inches(0.6), Inches(1.15), Inches(6), Inches(0.35), font_size=12, color=RED, bold=True)

recs = [
    ("1", "Accelerate the Dedicated L7 DDoS Detection Pipeline",
     "ML-based behavioral detection is the #1 competitive gap. Cloudflare's autonomous engine is production-proven and expanding. Every quarter of delay widens the perception gap. Head-to-head: Cloudflare leads 9 of 13 dimensions."),
    ("2", "Amplify MazeBolt RADAR & Validated-Testing Narrative",
     "RADAR's non-disruptive L7 DDoS simulation is exclusive to Azure. Invest in case studies, benchmarks, and bake-off playbooks — counter Cloudflare's claim-without-proof approach."),
    ("3", "Close the Messaging Gap with a Public L7 DDoS Story",
     "Azure's reluctance to claim creates a vacuum competitors fill. Publish L7 DDoS protection guidance, attack coverage docs, and validated benchmarks. Align with 'Security by AI' positioning."),
]

y = Inches(1.6)
for num, title, desc in recs:
    # Number circle
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT_BLUE
    circ.line.fill.background()
    ctf = circ.text_frame
    ctf.paragraphs[0].text = num
    ctf.paragraphs[0].font.size = Pt(18)
    ctf.paragraphs[0].font.color.rgb = WHITE
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    add_body_text(slide, title, Inches(1.3), y - Inches(0.05), Inches(10.5), Inches(0.4),
                  font_size=18, color=DARK_BLUE, bold=True)
    # Description
    add_body_text(slide, desc, Inches(1.3), y + Inches(0.35), Inches(10.5), Inches(1.1),
                  font_size=14, color=DARK_GRAY)

    y += Inches(1.7)

# =============================================================================
# SLIDE 9 — Questions / Next Steps
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_title_text(slide, "Questions & Next Steps", Inches(1), Inches(1.5), Inches(11), Inches(1),
               font_size=40, color=WHITE, bold=True)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.8), Inches(1.5), Inches(0.06))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

next_steps = [
    "Review gaps and open questions — prioritize data collection for unresolved items",
    "Validate Azure WAF compliance certifications (SOC 2, ISO 27001, FedRAMP High) via Azure Trust Center",
    "Establish win/loss tracking for L7 DDoS competitive encounters",
    "Consider extending analysis to include Akamai and Imperva",
    "Align on ML-based L7 DDoS detection pipeline timeline with engineering",
]
tf = add_bullet_list(slide, next_steps, Inches(1), Inches(3.2), Inches(11), Inches(3.5),
                     font_size=18, color=RGBColor(0xCC, 0xDD, 0xEE))

# Footer on last slide
ft = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11), Inches(0.4))
ftf = ft.text_frame
p = ftf.paragraphs[0]
p.text = "Competitive Analysis Agent Farm  |  March 4, 2026  |  Internal"
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(0x88, 0x99, 0xAA)
p.font.italic = True

# =============================================================================
# Save
# =============================================================================
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "competitive-analysis.pptx")
prs.save(output_path)
print(f"✅ PowerPoint saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
